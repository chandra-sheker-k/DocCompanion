
"""
PowerPoint (.pptx) document extractor.

Each slide is treated as a separate page for downstream
chunking and citation.

Returns:
    ExtractionResult
"""

from pathlib import Path
from typing import List

from pptx import Presentation
from .models import ExtractionResult, PageContent
from .utils import build_result

def extract(path: str) -> ExtractionResult:
    """
    Extract all text from a PowerPoint presentation.

    Parameters
    ----------
    path : str
        Absolute path to pptx file.

    Returns
    -------
    ExtractionResult
    """

    presentation = Presentation(path)
    pages: List[PageContent] = []
    warnings = []
    full_text = []
    slide_number = 1

    for slide in presentation.slides:
        slide_lines = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                slide_lines.append(text)
        slide_text = "\n".join(slide_lines)

        pages.append(
            PageContent(page=slide_number, text=slide_text)
        )

        full_text.append(slide_text)
        slide_number += 1
    metadata = {
        "filename": Path(path).name,
        "extension": ".pptx",
        "slides": len(presentation.slides),
        "pages": len(presentation.slides),
        "title": presentation.core_properties.title or "",
        "author": presentation.core_properties.author or "",
        "subject": presentation.core_properties.subject or "",
        "keywords": presentation.core_properties.keywords or "",
        "comments": presentation.core_properties.comments or "",
        "category": presentation.core_properties.category or "",
        "language": "",
        "created": str(presentation.core_properties.created)
        if presentation.core_properties.created
        else "",

        "modified": str(presentation.core_properties.modified)
        if presentation.core_properties.modified
        else ""
    }

    return build_result(text="\n\n".join(full_text), pages=pages,
        metadata=metadata, warnings=warnings)